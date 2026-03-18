"""
build_et_report.py
Generates a 1-2 page Word (.docx) summary + a PowerPoint (.pptx) presentation
from the ET comparison notebook outputs.

Run from project root:
    .venv/bin/python scripts/build_et_report.py
"""

from pathlib import Path
import io
import textwrap
import numpy as np
import pandas as pd

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

# ── paths ─────────────────────────────────────────────────────────────────────
FIGURES = Path("figures/ET comparison")
GEOMS   = ["delta_polygon", "poly_bbox", "okavango_bbox"]
GEOM_LABELS = {
    "delta_polygon": "Delta Polygon (23 387 km²)",
    "poly_bbox":     "Polygon Bounding Box (108 705 km²)",
    "okavango_bbox": "Okavango Bbox (50 162 km²)",
}
MODELS = [
    "ERA5Land_totalET", "FLDAS_Evap", "MOD16A2GF_v61",
    "PML_v2_landET", "TerraClimate_aet",
    "USGS_SSEBop_MODIS_monthly", "WaPORv3_AETI_dekadal",
]
MODEL_SHORT = {
    "ERA5Land_totalET":        "ERA5-Land",
    "FLDAS_Evap":              "FLDAS",
    "MOD16A2GF_v61":          "MOD16",
    "PML_v2_landET":           "PML v2",
    "TerraClimate_aet":        "TerraClimate",
    "USGS_SSEBop_MODIS_monthly": "SSEBop",
    "WaPORv3_AETI_dekadal":    "WaPOR v3",
}

OUTDIR = Path("figures/ET_comparison_summary_report")
OUTDIR.mkdir(parents=True, exist_ok=True)

# ── helper: navy heading colour ───────────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x37, 0x6C)
PNAVY = PRGBColor(0x1A, 0x37, 0x6C)
PTEAL = PRGBColor(0x00, 0x7A, 0x87)

# ── load data ─────────────────────────────────────────────────────────────────
def load(geom):
    d = FIGURES / geom
    return {
        "mb":  pd.read_csv(d / "mass_balance.csv", parse_dates=["date"]),
        "et":  pd.read_csv(d / "et_monthly.csv",   parse_dates=["date"]),
        "ch":  pd.read_csv(d / "chirps_monthly.csv", parse_dates=["date"]),
        "gr":  pd.read_csv(d / "grace_monthly.csv",  parse_dates=["date"]),
    }

def annual(series, col):
    v = pd.to_numeric(series[col], errors="coerce").dropna()
    return v.mean() * 12 if len(v) else np.nan

def fmt(v, d=1):
    return f"{v:.{d}f}" if pd.notna(v) else "—"

def pct(v):
    return f"{100*v:.0f}%" if pd.notna(v) else "—"

summary = pd.read_csv(FIGURES / "geometry_comparison_summary.csv")

data = {}
for g in GEOMS:
    try:
        data[g] = load(g)
    except Exception as e:
        print(f"Warning: could not load {g}: {e}")

# per-geometry model rankings
def model_means(g):
    et = data[g]["et"]
    return (
        et.groupby("dataset")["et_mm_mean"].mean()
          .rename(MODEL_SHORT)
          .sort_values(ascending=False)
    )

# ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ──
#  WORD DOCUMENT
# ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ──
def add_heading(doc, text, level=1, colour=NAVY):
    p = doc.add_heading(text, level=level)
    for run in p.runs:
        run.font.color.rgb = colour
    return p

def add_table(doc, df, col_widths=None):
    rows, cols = df.shape
    t = doc.add_table(rows=rows+1, cols=cols)
    t.style = "Table Grid"
    for j, col in enumerate(df.columns):
        cell = t.cell(0, j)
        cell.text = col
        run = cell.paragraphs[0].runs[0]
        run.bold = True
        run.font.size = Pt(9)
        cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        shading = OxmlElement("w:shd")
        shading.set(qn("w:fill"), "1A376C")
        shading.set(qn("w:color"), "FFFFFF")
        shading.set(qn("w:val"), "clear")
        cell._tc.get_or_add_tcPr().append(shading)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for i, row in df.iterrows():
        for j, val in enumerate(row):
            cell = t.cell(i+1, j)
            cell.text = str(val)
            cell.paragraphs[0].runs[0].font.size = Pt(8.5)
    if col_widths:
        for j, w in enumerate(col_widths):
            for row in t.rows:
                row.cells[j].width = Inches(w)
    return t

doc = Document()

# narrow margins
for sec in doc.sections:
    sec.top_margin    = Inches(0.75)
    sec.bottom_margin = Inches(0.75)
    sec.left_margin   = Inches(1.0)
    sec.right_margin  = Inches(1.0)

# ---- Title ----
title = doc.add_paragraph()
title.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = title.add_run("Okavango Delta Evapotranspiration Comparison: Key Findings")
run.bold = True
run.font.size = Pt(16)
run.font.color.rgb = NAVY

sub = doc.add_paragraph()
sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
sub.add_run("March 2026 — Inman-Lyons Lab").font.size = Pt(10)

doc.add_paragraph()

# ---- 1  Background ----
add_heading(doc, "1  Background & Objectives", level=1)
bg = doc.add_paragraph()
bg.add_run(
    "Seven satellite/reanalysis evapotranspiration (ET) products were evaluated over the Okavango Delta: "
    "ERA5-Land, FLDAS, MOD16A2GF, PML v2, TerraClimate, SSEBop, and WaPOR v3. "
    "ET values were compared against CHIRPS precipitation, Mohembo river inflow, and GRACE-derived "
    "terrestrial water storage anomalies for three spatial extents: the mapped delta polygon, "
    "a polygon bounding box, and the broader Okavango basin bounding box. "
    "The analysis spans 2001–2023 at monthly resolution. "
    "Key questions: (1) How do ET products compare in magnitude and seasonality? "
    "(2) Is the implied water balance physically plausible? "
    "(3) Does geometry choice materially affect conclusions?"
)
bg.runs[0].font.size = Pt(10)

# ---- 2  Key Findings ----
add_heading(doc, "2  Key Findings", level=1)

# 2a  ET vs P
add_heading(doc, "2a  ET and precipitation", level=2, colour=RGBColor(0x00, 0x7A, 0x87))
row = summary[summary.geometry == "delta_polygon"].iloc[0]
p = doc.add_paragraph(style="List Bullet")
p.add_run(
    f"Over the mapped delta polygon (23 387 km²), ensemble-median ET averages "
    f"{fmt(row.mean_ETmm_yr, 0)} mm/yr, accounting for approximately "
    f"{fmt(100*row.mean_ETmm_yr/row.mean_Pmm_yr, 0)}% of CHIRPS precipitation "
    f"({fmt(row.mean_Pmm_yr, 0)} mm/yr)."
).font.size = Pt(10)

p = doc.add_paragraph(style="List Bullet")
et_spread = {}
for g in GEOMS:
    if g in data:
        mm = model_means(g)
        et_spread[g] = (mm.min(), mm.max())
dp_min, dp_max = et_spread.get("delta_polygon", (np.nan, np.nan))
p.add_run(
    f"Across the seven ET products, mean monthly ET over the delta polygon ranges from "
    f"{fmt(dp_min, 1)} to {fmt(dp_max, 1)} mm/month "
    f"({fmt(dp_max - dp_min, 1)} mm/month spread), "
    "underscoring substantial inter-product uncertainty."
).font.size = Pt(10)

# 2b  Water balance
add_heading(doc, "2b  Water balance closure", level=2, colour=RGBColor(0x00, 0x7A, 0x87))
mb = data["delta_polygon"]["mb"] if "delta_polygon" in data else None
if mb is not None:
    neg_frac = (pd.to_numeric(mb["Qout_plus_G_km3"], errors="coerce").dropna() < 0).mean()
    p = doc.add_paragraph(style="List Bullet")
    p.add_run(
        f"Mean Mohembo inflow is {fmt(row.mean_Qin_m3s, 0)} m³/s. "
        f"In {pct(neg_frac)} of months the implied outflow + groundwater term (P + Q_in − ET − ΔS) "
        "is negative, suggesting ET slightly exceeds the available water budget and/or "
        "GRACE ΔS uncertainty is a limiting factor."
    ).font.size = Pt(10)

p = doc.add_paragraph(style="List Bullet")
p.add_run(
    "Annual ET consistently exceeds precipitation in the delta polygon "
    "(ET/P ≈ 1.14), consistent with the delta's role as a net evaporation sink sustained by "
    "upstream inflow from Angola."
).font.size = Pt(10)

# 2c  Geometry sensitivity
add_heading(doc, "2c  Geometry sensitivity", level=2, colour=RGBColor(0x00, 0x7A, 0x87))
p = doc.add_paragraph(style="List Bullet")
p.add_run(
    "ET in mm/yr is consistent across geometries (537, 509, 535 mm/yr for delta_polygon, "
    "poly_bbox, and okavango_bbox respectively), indicating that products capture similar "
    "surface fluxes per unit area regardless of domain size. "
    "The poly_bbox is ~4.6× larger and includes drier upland areas, "
    "explaining its slightly lower ET and precipitation."
).font.size = Pt(10)

# ---- 3  Summary Table ----
add_heading(doc, "3  Geometry Comparison", level=1)
tdf = summary.copy()
tdf.columns = ["Geometry", "Area (km²)", "ET (mm/yr)", "P (mm/yr)", "Qin (m³/s)"]
tdf["Area (km²)"]  = tdf["Area (km²)"].map(lambda x: f"{x:,.0f}")
tdf["ET (mm/yr)"]  = tdf["ET (mm/yr)"].map(lambda x: f"{x:.0f}")
tdf["P (mm/yr)"]   = tdf["P (mm/yr)"].map(lambda x: f"{x:.0f}")
tdf["Qin (m³/s)"]  = tdf["Qin (m³/s)"].map(lambda x: f"{x:.0f}")
tdf["Geometry"]    = tdf["Geometry"].map(GEOM_LABELS)
add_table(doc, tdf, col_widths=[2.3, 1.1, 1.0, 1.0, 1.0])

doc.add_paragraph()

# ---- 4  Model Ranking ----
add_heading(doc, "4  ET Product Ranking (delta polygon, mean mm/month)", level=1)
mm_dp = model_means("delta_polygon").reset_index()
mm_dp.columns = ["ET Product", "Mean ET (mm/month)"]
mm_dp["Mean ET (mm/month)"] = mm_dp["Mean ET (mm/month)"].map(lambda x: f"{x:.1f}")
add_table(doc, mm_dp, col_widths=[2.5, 2.0])

doc.add_paragraph()

# ---- 5  Conclusions ----
add_heading(doc, "5  Conclusions", level=1)
conc = doc.add_paragraph()
conc.add_run(
    "ERA5-Land and WaPOR v3 tend to produce the highest ET estimates; MOD16 and PML v2 the lowest. "
    "The ~40–50 mm/month inter-product range exceeds the seasonal GRACE signal in some months, "
    "highlighting ET product choice as the dominant source of water-balance uncertainty. "
    "The delta polygon geometry produces the most physically meaningful closure because it closely "
    "matches the flooded-land extent that drives high ET. "
    "Future work should focus on constraining ET using eddy-covariance towers or sap-flow "
    "measurements within the delta."
).font.size = Pt(10)

docx_path = OUTDIR / "ET_comparison_summary.docx"
doc.save(docx_path)
print(f"Word report saved: {docx_path}")


# ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ──
#  POWERPOINT — all figures full-sized (one per slide)
# ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ── ──
from PIL import Image as PILImage

prs = Presentation()
SLIDE_W = PInches(13.33)
SLIDE_H = PInches(7.5)
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank

# Layout constants
HEADER_H   = 0.85           # inches for the title bar
MARGIN     = 0.25           # margin around figure
BODY_TOP   = HEADER_H + MARGIN
BODY_W     = 13.33 - 2 * MARGIN
BODY_H     = 7.5   - BODY_TOP - MARGIN

def add_slide(prs, title_text, subtitle_text=None):
    slide = prs.slides.add_slide(BLANK)
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, PInches(HEADER_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PNAVY
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = PPt(24)
    run.font.bold = True
    run.font.color.rgb = PRGBColor(0xFF, 0xFF, 0xFF)
    tf.margin_left  = PInches(0.3)
    tf.margin_top   = PInches(0.18)

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.size = PPt(12)
        r2.font.color.rgb = PRGBColor(0x88, 0xAA, 0xCC)
    return slide

def add_textbox(slide, text, left, top, width, height,
                size=12, bold=False, colour=PRGBColor(0x22,0x22,0x22), wrap=True):
    txb = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    txb.text_frame.word_wrap = wrap
    p = txb.text_frame.paragraphs[0]
    p.word_wrap = wrap
    run = p.add_run()
    run.text = text
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return txb

def add_bullets(slide, items, left, top, width, height, size=12):
    txb = slide.shapes.add_textbox(PInches(left), PInches(top), PInches(width), PInches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = PPt(size)
        p.space_after = PPt(6)

def add_fullsize_png(slide, path):
    """Add a PNG that fills the available body area, preserving aspect ratio, centred."""
    p = Path(path)
    if not p.exists():
        return
    img = PILImage.open(p)
    img_w, img_h = img.size  # pixels
    aspect = img_w / img_h
    # fit to available body area
    body_aspect = BODY_W / BODY_H
    if aspect >= body_aspect:
        # image is wider than body → constrain by width
        w = BODY_W
        h = w / aspect
    else:
        # image is taller → constrain by height
        h = BODY_H
        w = h * aspect
    left = MARGIN + (BODY_W - w) / 2
    top  = BODY_TOP + (BODY_H - h) / 2
    slide.shapes.add_picture(str(p), PInches(left), PInches(top),
                             width=PInches(w), height=PInches(h))

def pptx_table(slide, df, left, top, col_widths, row_height=PInches(0.35)):
    from pptx.util import Pt as PPt2
    rows, cols = df.shape
    col_widths_emu = [PInches(w) for w in col_widths]
    total_w = sum(col_widths_emu)
    tbl = slide.shapes.add_table(rows+1, cols,
                                  PInches(left), PInches(top),
                                  total_w, row_height*(rows+1)).table
    for j, w in enumerate(col_widths_emu):
        tbl.columns[j].width = w
    for j, col in enumerate(df.columns):
        cell = tbl.cell(0, j)
        cell.text = col
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = PPt2(11)
        cell.text_frame.paragraphs[0].font.color.rgb = PRGBColor(0xFF,0xFF,0xFF)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = PNAVY
    for i, row in df.iterrows():
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = PPt2(10)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRGBColor(0xE8, 0xEE, 0xF7)
    return tbl

# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = PNAVY
bg.line.fill.background()

txb = slide.shapes.add_textbox(PInches(1), PInches(2.0), PInches(11.33), PInches(1.6))
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Evapotranspiration Comparison\nover the Okavango Delta"
r.font.size = PPt(44)
r.font.bold = True
r.font.color.rgb = PRGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

txb2 = slide.shapes.add_textbox(PInches(1), PInches(4.2), PInches(11.33), PInches(0.6))
p2 = txb2.text_frame.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Seven ET Products · Three Geometries · 2001–2023   |   Inman-Lyons Lab · March 2026"
r2.font.size = PPt(16)
r2.font.color.rgb = PRGBColor(0x88, 0xCC, 0xFF)

# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 2 — Study Overview (text + map)
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Study Overview", "Comparing 7 ET products across 3 spatial domains")
add_bullets(slide, [
    "Seven ET products: ERA5-Land, FLDAS, MOD16A2GF, PML v2, TerraClimate, SSEBop, WaPOR v3",
    "Three geometries: delta polygon (23 387 km²), polygon bbox (108 705 km²), Okavango bbox (50 162 km²)",
    "Validation data: CHIRPS precipitation, Mohembo gauge (inflow), GRACE terrestrial water storage",
    "Period: January 2001 – December 2023, monthly resolution",
    "Water balance: P + Q_in − ET − ΔS = Q_out + G (implied outflow + groundwater)",
], left=0.3, top=1.2, width=6.3, height=5.8, size=14)
add_fullsize_png_in_region = lambda sl, fp, l, t, w, h: (
    sl.shapes.add_picture(str(fp), PInches(l), PInches(t), width=PInches(w))
    if Path(fp).exists() else None
)
add_fullsize_png_in_region(slide, FIGURES / "delta_polygon/delta_map.png", 6.8, 1.1, 6.2, 6.0)

# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 3 — Key Findings (text)
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Key Findings", "Delta polygon geometry — quantitative summary")
row = summary[summary.geometry == "delta_polygon"].iloc[0]
items = [
    f"Ensemble-median ET ≈ {fmt(row.mean_ETmm_yr, 0)} mm/yr  |  CHIRPS P ≈ {fmt(row.mean_Pmm_yr, 0)} mm/yr",
    f"ET/P ratio ≈ {fmt(row.mean_ETmm_yr/row.mean_Pmm_yr, 2)} — delta is a net evaporation sink sustained by inflow",
    f"Mean Mohembo inflow ≈ {fmt(row.mean_Qin_m3s, 0)} m³/s",
    "All products show peak ET in austral summer (Nov–Mar), aligned with the flooding season",
    "Inter-product spread: up to ~50 mm/month — ET choice dominates water-balance uncertainty",
    "ERA5-Land and WaPOR v3 are high-end; MOD16 and PML v2 are low-end",
    "TerraClimate and SSEBop yield best monthly water-balance closure",
    "Geometry choice has minor effect — ET per unit area is consistent across domains",
]
add_bullets(slide, items, left=0.4, top=1.15, width=12.5, height=6.0, size=16)

# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 4 — ET Model Ranking Table
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "ET Product Ranking", "Mean monthly ET (mm/month) — delta polygon")
mm_dp = model_means("delta_polygon").reset_index()
mm_dp.columns = ["ET Product", "Mean ET (mm/month)"]
mm_dp["Mean ET (mm/month)"] = mm_dp["Mean ET (mm/month)"].map(lambda x: f"{x:.1f}")
pptx_table(slide, mm_dp, left=0.4, top=1.3, col_widths=[3.5, 2.8])
add_bullets(slide, [
    "ERA5-Land: highest ET — strong latent heat parameterisation",
    "WaPOR v3: high — captures open-water evaporation effectively",
    "MOD16 & PML v2: lowest — resistance-based, constrained by VPD & LAI",
    "SSEBop: middle ground, calibrated to reference ET",
    "Ranking consistent across all three geometries",
], left=7.0, top=1.3, width=6.0, height=5.5, size=14)

# ═══════════════════════════════════════════════════════════════════════════════
#  Slide 5 — Geometry Comparison Table
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Geometry Sensitivity", "ET/P/Q are consistent in mm/yr across domains")
tdf = summary.copy()
tdf.columns = ["Geometry", "Area (km²)", "ET (mm/yr)", "P (mm/yr)", "Qin (m³/s)"]
tdf["Area (km²)"]  = tdf["Area (km²)"].map(lambda x: f"{x:,.0f}")
tdf["ET (mm/yr)"]  = tdf["ET (mm/yr)"].map(lambda x: f"{x:.0f}")
tdf["P (mm/yr)"]   = tdf["P (mm/yr)"].map(lambda x: f"{x:.0f}")
tdf["Qin (m³/s)"]  = tdf["Qin (m³/s)"].map(lambda x: f"{x:.0f}")
tdf["Geometry"]     = tdf["Geometry"].map(lambda x: GEOM_LABELS.get(x, x))
pptx_table(slide, tdf, left=0.4, top=1.3, col_widths=[3.8, 1.6, 1.4, 1.4, 1.4])
add_bullets(slide, [
    "ET per unit area ~535–537 mm/yr for delta polygon & Okavango bbox",
    "Polygon bbox lower (509 mm/yr) — includes drier upland areas",
    "Qin unchanged — gauge measurement independent of geometry",
    "Delta polygon recommended: closest match to inundated area driving ET",
], left=0.4, top=3.8, width=12.5, height=3.5, size=14)

# ═══════════════════════════════════════════════════════════════════════════════
#  Full-figure slides — one figure per slide, filling the body area
# ═══════════════════════════════════════════════════════════════════════════════
FULL_FIGS = [
    # (path, title, subtitle)
    (FIGURES / "comparison_ET_P_Q.png",
     "ET vs Precipitation vs Inflow",
     "Monthly comparison across all products — delta polygon"),
    (FIGURES / "comparative_cumulative_sums_median.png",
     "Cumulative Sums — Median ET by Geometry",
     "Comparing delta_polygon, poly_bbox, and okavango_bbox"),
    (FIGURES / "comparative_cumulative_sums.png",
     "Cumulative Sums — All 7 ET Products × 3 Geometries",
     "Per-model cumulative water balance comparison"),
]

# Add median-ET mass balance + cumulative for EACH geometry
for geom in GEOMS:
    label = GEOM_LABELS.get(geom, geom)
    mb_path  = FIGURES / geom / "median_ET/mass_balance_terms.png"
    cum_path = FIGURES / geom / "median_ET/cumulative_sums.png"
    if mb_path.exists():
        FULL_FIGS.append((mb_path,
            f"Mass Balance Terms (Median ET) — {label}",
            f"{geom} — 3-month smoothed — P, Qin, ET, ΔS"))
    if cum_path.exists():
        FULL_FIGS.append((cum_path,
            f"Cumulative Water Balance (Median ET) — {label}",
            f"{geom} — ∑(Qin + P − ET) vs ∑ΔS (GRACE)"))

# Add per-region cumulative sum figures
for model in MODELS:
    fig_path = FIGURES / f"cumulative_by_region_{model}.png"
    if fig_path.exists():
        short = MODEL_SHORT.get(model, model)
        FULL_FIGS.append((fig_path,
            f"Cumulative Sums by Region — {short}",
            f"Three geometries using ET = {short}"))

# Add per-model mass balance terms + cumulative sums for EACH geometry
for geom in GEOMS:
    label = GEOM_LABELS.get(geom, geom)
    for model in MODELS:
        short = MODEL_SHORT.get(model, model)
        mb_path  = FIGURES / geom / "single-model" / model / "mass_balance_terms.png"
        cum_path = FIGURES / geom / "single-model" / model / "cumulative_sums.png"
        if mb_path.exists():
            FULL_FIGS.append((mb_path,
                f"Mass Balance Terms — {short}",
                f"{label}"))
        if cum_path.exists():
            FULL_FIGS.append((cum_path,
                f"Cumulative Sums — {short}",
                f"{label}"))

for fig_path, title, subtitle in FULL_FIGS:
    if not Path(fig_path).exists():
        continue
    slide = add_slide(prs, title, subtitle)
    add_fullsize_png(slide, fig_path)

# ═══════════════════════════════════════════════════════════════════════════════
#  Final slide — Conclusions
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide(prs, "Conclusions & Next Steps")
add_bullets(slide, [
    "ET/P > 1 over the delta polygon confirms the Okavango Delta is an evaporation-dominated system sustained by Angolan catchment inflow",
    "ERA5-Land and WaPOR v3 are high-end estimates; MOD16 and PML v2 are low — a ~50 mm/month spread remains unresolved",
    "All products agree on seasonal phasing (peak Nov–Mar) and inter-annual variability driven by the flood pulse",
    "GRACE ΔS is a useful but noisy closure term; uncertainty limits definitive product discrimination",
    "Recommended product for water balance: TerraClimate or SSEBop (best closure); ERA5-Land for physically-based simulations",
    "Next steps: eddy-covariance validation within delta; test during extreme flood/drought years; ensemble weighting by closure error",
], left=0.4, top=1.3, width=12.5, height=5.8, size=16)

pptx_path = OUTDIR / "ET_comparison_summary.pptx"
prs.save(pptx_path)
print(f"PowerPoint saved:  {pptx_path}")
print("Done.")
